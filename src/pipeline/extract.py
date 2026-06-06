"""File extraction: PDF, DOCX, TXT → raw text with structural cleanup."""

import os
import tempfile
from pathlib import Path
from typing import Optional

from docx import Document as DocxDocument
from pypdf import PdfReader


class ExtractionError(Exception):
    pass


def extract_text(file_path: str | Path, mime_type: Optional[str] = None) -> str:
    """Extract text from PDF, DOCX, or TXT with structural cleanup."""
    path = Path(file_path)
    if not path.exists():
        raise ExtractionError(f"File not found: {file_path}")

    # Guard against very large files that could OOM the container (fix #12)
    max_bytes = 10 * 1024 * 1024  # 10 MB
    file_size = path.stat().st_size
    if file_size > max_bytes:
        raise ExtractionError(
            f"File '{path.name}' is {file_size // (1024*1024):.1f} MB — "
            f"maximum supported size is 10 MB. Please split the document and re-upload."
        )

    inferred = mime_type or _infer_type(path)
    text = ""

    try:
        if inferred == "application/pdf":
            text = _extract_pdf(path)
        elif inferred in (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/msword",
        ):
            text = _extract_docx(path)
        elif inferred == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            text = _extract_pptx(path)
        elif inferred == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            text = _extract_xlsx(path)
        elif inferred == "text/plain":
            text = _extract_txt(path)
        else:
            raise ExtractionError(f"Unsupported file type: {inferred}")
    except ExtractionError:
        raise
    except Exception as e:
        raise ExtractionError(f"Failed to extract {path.name}: {e}")

    return text.strip()


def _infer_type(path: Path) -> str:
    ext = path.suffix.lower()
    if ext == ".pdf":
        return "application/pdf"
    elif ext in (".docx", ".doc"):
        return "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    elif ext in (".pptx", ".ppt"):
        return "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    elif ext in (".xlsx", ".xls"):
        return "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    elif ext == ".txt":
        return "text/plain"
    elif ext in (".md", ".markdown"):
        return "text/plain"
    return "application/octet-stream"


def _extract_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    lines = []
    for page in reader.pages:
        page_text = page.extract_text()
        if not page_text:
            continue
        for line in page_text.split("\n"):
            cleaned = line.strip()
            if cleaned:
                lines.append(cleaned)
        lines.append("")
    return "\n".join(lines)


_TOC_STYLES = {"toc 2", "toc 3", "toc 4", "toc 5", "toc 6"}


def _extract_docx(path: Path) -> str:
    from docx.text.paragraph import Paragraph as DocxParagraph
    from docx.table import Table as DocxTable

    doc = DocxDocument(str(path))
    parts = []

    for child in doc.element.body.iterchildren():
        tag = child.tag.split("}")[-1]

        if tag == "p":
            para = DocxParagraph(child, doc)
            style = para.style.name if para.style else ""
            if style in _TOC_STYLES:
                continue
            text = para.text.strip()
            if not text:
                continue
            if style.startswith("Heading"):
                try:
                    level = int(style.split()[-1])
                    prefix = "#" * min(level, 4) + " "
                except (ValueError, IndexError):
                    prefix = "## "
                parts.append(prefix + text)
            elif "List" in style:
                parts.append("- " + text)
            else:
                parts.append(text)

        elif tag == "tbl":
            table = DocxTable(child, doc)
            table_lines = []
            for i, row in enumerate(table.rows):
                seen = set()
                cells = []
                for cell in row.cells:
                    cid = id(cell._tc)
                    if cid not in seen:
                        seen.add(cid)
                        # Clean cell text and replace inner newlines with spaces or HTML breaks
                        cell_txt = cell.text.strip().replace("\n", " ").replace("|", "\\|")
                        cells.append(cell_txt)
                if any(cells):
                    table_lines.append("| " + " | ".join(cells) + " |")
                    if i == 0:
                        # Add Markdown table divider row
                        table_lines.append("| " + " | ".join(["---"] * len(cells)) + " |")
            if table_lines:
                parts.extend(table_lines)
                parts.append("")  # Empty line after the table

    return "\n".join(parts)


def _extract_txt(path: Path) -> str:
    encodings = ["utf-8", "latin-1", "cp1252"]
    for enc in encodings:
        try:
            return open(path, "r", encoding=enc).read()
        except UnicodeDecodeError:
            continue
    raise ExtractionError(f"Could not decode {path.name} with any known encoding")


def chunk_text(text: str, chunk_size: int = 800, overlap: int = 100) -> list[dict]:
    """Split text into overlapping chunks with context headers."""
    if not text.strip():
        return []

    lines = text.split("\n")
    chunks = []
    start = 0
    chunk_num = 0

    while start < len(text):
        end = start + chunk_size
        chunk_text = text[start:end]

        first_line = chunk_text.split("\n")[0][:100]
        chunks.append(
            {
                "chunk_id": f"chunk-{chunk_num:04d}",
                "text": chunk_text,
                "char_start": start,
                "char_end": end,
                "preview": first_line,
            }
        )

        start += chunk_size - overlap
        chunk_num += 1

    return chunks


def save_upload(upload_file, destination: Path) -> Path:
    """Save an uploaded file to a destination path."""
    with open(destination, "wb") as f:
        f.write(upload_file.read())
    return destination


def _extract_pptx(path: Path) -> str:
    """Extract text from PowerPoint presentations page-by-page."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ExtractionError("python-pptx is not installed. Run: pip install python-pptx")

    prs = Presentation(str(path))
    slides_text = []
    for i, slide in enumerate(prs.slides):
        slide_lines = [f"--- Slide {i+1} ---"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_lines.append(shape.text.strip())
        slides_text.append("\n".join(slide_lines))
    return "\n\n".join(slides_text)


def _extract_xlsx(path: Path) -> str:
    """Extract text from Excel sheets, formatting tables as markdown grids."""
    try:
        import openpyxl
    except ImportError:
        raise ExtractionError("openpyxl is not installed. Run: pip install openpyxl")

    wb = openpyxl.load_workbook(filename=str(path), read_only=True, data_only=True)
    sheets_text = []
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        sheets_text.append(f"### Sheet: {sheet_name}")
        
        # Read rows and filter empty ones
        rows_data = []
        for row in sheet.iter_rows(values_only=True):
            row_vals = [str(v).strip() if v is not None else "" for v in row]
            if any(row_vals):
                rows_data.append(row_vals)
        
        # Convert to Markdown table format
        if rows_data:
            table_lines = []
            for i, row in enumerate(rows_data):
                table_lines.append("| " + " | ".join(row).replace("\n", " ") + " |")
                if i == 0:
                    table_lines.append("| " + " | ".join(["---"] * len(row)) + " |")
            sheets_text.extend(table_lines)
            sheets_text.append("")
    return "\n".join(sheets_text)
